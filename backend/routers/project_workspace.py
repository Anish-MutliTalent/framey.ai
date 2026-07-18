"""
Project workspace: files, notes, activity log, AI research, file conversion.
Files are stored at  backend/uploads/projects/{project_id}/
"""
import io
import json
import mimetypes
import os
import uuid
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, HTMLResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session

import models
from auth import get_current_user, get_file_user
from database import get_db

router = APIRouter(tags=["project-workspace"])

UPLOAD_ROOT = Path(__file__).parent.parent / "uploads" / "projects"
UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)

MAX_FILE_MB = 50
ALLOWED_TYPES = {
    ".pdf":  "pdf",
    ".html": "html",
    ".htm":  "html",
    ".pptx": "pptx",
    ".ppt":  "pptx",
    ".docx": "docx",
    ".doc":  "docx",
    ".png":  "image",
    ".jpg":  "image",
    ".jpeg": "image",
    ".gif":  "image",
    ".svg":  "image",
    ".txt":  "text",
    ".md":   "text",
    ".csv":  "text",
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _project_dir(project_id: int) -> Path:
    d = UPLOAD_ROOT / str(project_id)
    d.mkdir(parents=True, exist_ok=True)
    return d


def _log(db: Session, project_id: int, user_id: int, action: str, details: str = ""):
    db.add(models.ProjectActivity(
        project_id=project_id,
        user_id=user_id,
        action=action,
        details=details,
    ))


def _fmt_file(f: models.ProjectFile) -> dict:
    return {
        "id": f.id,
        "original_name": f.original_name,
        "file_type": f.file_type,
        "mime_type": f.mime_type,
        "size_bytes": f.size_bytes,
        "is_submitted": f.is_submitted,
        "folder": f.folder,
        "uploaded_by": f.uploader.name if f.uploader else "",
        "created_at": f.created_at.isoformat(),
    }


def _fmt_note(n: models.ProjectNote) -> dict:
    return {
        "id": n.id,
        "title": n.title,
        "content": n.content,
        "created_at": n.created_at.isoformat(),
        "updated_at": n.updated_at.isoformat() if n.updated_at else n.created_at.isoformat(),
    }


def _assert_project_access(project_id: int, user: models.User, db: Session) -> models.Project:
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if not project:
        raise HTTPException(404, "Project not found")
    user_roles = {r.role for r in user.roles}
    is_owner = project.user_id == user.id
    is_teacher = user_roles.intersection({"teacher", "class_teacher", "coordinator", "principal", "tech_admin"})
    if not is_owner and not is_teacher:
        raise HTTPException(403, "Access denied")
    return project


# ════════════════════════════════════════════════════
#  FILE UPLOAD & LISTING
# ════════════════════════════════════════════════════

@router.post("/projects/{project_id}/files")
async def upload_file(
    project_id: int,
    file: UploadFile = File(...),
    folder: str = Form(default=""),
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _assert_project_access(project_id, current_user, db)

    suffix = Path(file.filename or "file").suffix.lower()
    file_type = ALLOWED_TYPES.get(suffix, "other")

    content = await file.read()
    if len(content) > MAX_FILE_MB * 1024 * 1024:
        raise HTTPException(413, f"File exceeds {MAX_FILE_MB}MB limit")

    stored_name = f"{uuid.uuid4().hex}{suffix}"
    dest = _project_dir(project_id) / stored_name
    dest.write_bytes(content)

    mime = mimetypes.guess_type(file.filename or "")[0] or "application/octet-stream"

    pf = models.ProjectFile(
        project_id=project_id,
        uploaded_by=current_user.id,
        original_name=file.filename or stored_name,
        stored_name=stored_name,
        file_type=file_type,
        mime_type=mime,
        size_bytes=len(content),
        folder=folder,
    )
    db.add(pf)
    db.flush()
    _log(db, project_id, current_user.id, "uploaded file", pf.original_name)
    db.commit()
    db.refresh(pf)
    return _fmt_file(pf)


@router.get("/projects/{project_id}/files")
def list_files(
    project_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _assert_project_access(project_id, current_user, db)
    files = (
        db.query(models.ProjectFile)
        .filter(models.ProjectFile.project_id == project_id)
        .order_by(models.ProjectFile.folder, models.ProjectFile.created_at)
        .all()
    )
    return [_fmt_file(f) for f in files]


@router.delete("/project-files/{file_id}", status_code=204)
def delete_file(
    file_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    pf = db.query(models.ProjectFile).filter(models.ProjectFile.id == file_id).first()
    if not pf:
        raise HTTPException(404, "File not found")
    _assert_project_access(pf.project_id, current_user, db)

    path = _project_dir(pf.project_id) / pf.stored_name
    if path.exists():
        path.unlink()

    _log(db, pf.project_id, current_user.id, "deleted file", pf.original_name)
    db.delete(pf)
    db.commit()


@router.patch("/project-files/{file_id}/submit")
def toggle_submit(
    file_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    pf = db.query(models.ProjectFile).filter(models.ProjectFile.id == file_id).first()
    if not pf:
        raise HTTPException(404, "File not found")
    _assert_project_access(pf.project_id, current_user, db)
    pf.is_submitted = not pf.is_submitted
    action = "submitted file" if pf.is_submitted else "unsubmitted file"
    _log(db, pf.project_id, current_user.id, action, pf.original_name)
    db.commit()
    return {"is_submitted": pf.is_submitted}


# ════════════════════════════════════════════════════
#  FILE SERVING  (raw download or view)
# ════════════════════════════════════════════════════

@router.get("/project-files/{file_id}/raw")
def serve_file(
    file_id: int,
    current_user: models.User = Depends(get_file_user),   # accepts ?token= for browser opens
    db: Session = Depends(get_db),
):
    pf = db.query(models.ProjectFile).filter(models.ProjectFile.id == file_id).first()
    if not pf:
        raise HTTPException(404, "File not found")
    _assert_project_access(pf.project_id, current_user, db)
    path = _project_dir(pf.project_id) / pf.stored_name
    if not path.exists():
        raise HTTPException(404, "File not on disk")
    _log(db, pf.project_id, current_user.id, "viewed file", pf.original_name)
    db.commit()
    return FileResponse(str(path), media_type=pf.mime_type or "application/octet-stream",
                        filename=pf.original_name)


# ════════════════════════════════════════════════════
#  FILE CONVERSION  (PPTX → slides JSON, DOCX → HTML)
# ════════════════════════════════════════════════════

@router.get("/project-files/{file_id}/convert")
def convert_file(
    file_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """Convert PPTX to slide JSON or DOCX to HTML. Returns {type, data}."""
    pf = db.query(models.ProjectFile).filter(models.ProjectFile.id == file_id).first()
    if not pf:
        raise HTTPException(404, "File not found")
    _assert_project_access(pf.project_id, current_user, db)
    path = _project_dir(pf.project_id) / pf.stored_name
    if not path.exists():
        raise HTTPException(404, "File not on disk")

    if pf.file_type == "pptx":
        return {"type": "pptx", "slides": _pptx_to_slides(path)}
    elif pf.file_type == "docx":
        return {"type": "docx", "html": _docx_to_html(path)}
    elif pf.file_type == "text":
        return {"type": "text", "content": path.read_text(errors="replace")}
    else:
        raise HTTPException(400, f"Conversion not supported for type '{pf.file_type}'")


def _pptx_to_slides(path: Path) -> list:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        prs = Presentation(str(path))
    except Exception as e:
        return [{"title": "Error", "body": [str(e)], "notes": ""}]

    slides = []
    for i, slide in enumerate(prs.slides):
        title_text = ""
        body_lines = []
        notes_text = ""

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            # Heuristic: title placeholder or first shape
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph and ph.idx == 0:   # idx 0 = title
                title_text = text
            else:
                # Split by newlines and add as bullets
                for line in text.split("\n"):
                    line = line.strip()
                    if line:
                        body_lines.append(line)

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""

        slides.append({
            "index": i + 1,
            "title": title_text or f"Slide {i + 1}",
            "body": body_lines,
            "notes": notes_text,
        })
    return slides


def _docx_to_html(path: Path) -> str:
    try:
        import mammoth
        with open(str(path), "rb") as f:
            result = mammoth.convert_to_html(f)
        return result.value
    except Exception as e:
        return f"<p style='color:red'>Error converting document: {e}</p>"


# ════════════════════════════════════════════════════
#  PROJECT NOTES  (rich text, per project)
# ════════════════════════════════════════════════════

class NoteCreate(BaseModel):
    title: str = "Untitled Note"
    content: str = ""

class NoteUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None


@router.get("/projects/{project_id}/notes")
def list_notes(
    project_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _assert_project_access(project_id, current_user, db)
    notes = (
        db.query(models.ProjectNote)
        .filter(models.ProjectNote.project_id == project_id)
        .order_by(models.ProjectNote.updated_at.desc())
        .all()
    )
    return [_fmt_note(n) for n in notes]


@router.post("/projects/{project_id}/notes")
def create_note(
    project_id: int,
    payload: NoteCreate,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _assert_project_access(project_id, current_user, db)
    note = models.ProjectNote(
        project_id=project_id,
        author_id=current_user.id,
        title=payload.title,
        content=payload.content,
    )
    db.add(note)
    db.flush()
    _log(db, project_id, current_user.id, "created note", payload.title)
    db.commit()
    db.refresh(note)
    return _fmt_note(note)


@router.put("/project-notes/{note_id}")
def update_note(
    note_id: int,
    payload: NoteUpdate,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    note = db.query(models.ProjectNote).filter(models.ProjectNote.id == note_id).first()
    if not note:
        raise HTTPException(404, "Note not found")
    _assert_project_access(note.project_id, current_user, db)
    if payload.title is not None:
        note.title = payload.title
    if payload.content is not None:
        note.content = payload.content
    db.commit()
    db.refresh(note)
    return _fmt_note(note)


@router.delete("/project-notes/{note_id}", status_code=204)
def delete_note(
    note_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    note = db.query(models.ProjectNote).filter(models.ProjectNote.id == note_id).first()
    if not note:
        raise HTTPException(404, "Note not found")
    _assert_project_access(note.project_id, current_user, db)
    _log(db, note.project_id, current_user.id, "deleted note", note.title)
    db.delete(note)
    db.commit()


# ════════════════════════════════════════════════════
#  ACTIVITY LOG
# ════════════════════════════════════════════════════

@router.get("/projects/{project_id}/activity")
def get_activity(
    project_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _assert_project_access(project_id, current_user, db)
    logs = (
        db.query(models.ProjectActivity)
        .filter(models.ProjectActivity.project_id == project_id)
        .order_by(models.ProjectActivity.created_at.desc())
        .limit(100)
        .all()
    )
    return [
        {
            "id": l.id,
            "user": l.user.name,
            "action": l.action,
            "details": l.details,
            "created_at": l.created_at.isoformat(),
        }
        for l in logs
    ]


# ════════════════════════════════════════════════════
#  AI RESEARCH (scoped to project)
# ════════════════════════════════════════════════════

class ResearchMessage(BaseModel):
    message: str
    history: list = []


@router.post("/projects/{project_id}/ai-research")
def ai_research(
    project_id: int,
    payload: ResearchMessage,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    import anthropic as anthropic_lib
    project = _assert_project_access(project_id, current_user, db)

    # Gather context: file names, notes titles
    files = db.query(models.ProjectFile).filter(models.ProjectFile.project_id == project_id).all()
    notes = db.query(models.ProjectNote).filter(models.ProjectNote.project_id == project_id).all()
    file_list = ", ".join(f.original_name for f in files) or "none"
    notes_list = ", ".join(n.title for n in notes) or "none"

    system = (
        f"You are a research assistant helping with the project: '{project.title}' "
        f"(Subject: {project.subject or 'general'}, Description: {project.description or 'not set'}).\n"
        f"Files in project: {file_list}\n"
        f"Notes: {notes_list}\n\n"
        "Help the student research, brainstorm, write, and analyse. Be concise and academic."
    )

    client = anthropic_lib.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", ""))
    messages = [
        {"role": "user" if m.get("role") == "user" else "assistant", "content": m.get("content", "")}
        for m in payload.history
    ]
    messages.append({"role": "user", "content": payload.message})

    try:
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1000,
            system=system,
            messages=messages,
        )
        _log(db, project_id, current_user.id, "used AI research", payload.message[:80])
        db.commit()
        return {"response": resp.content[0].text}
    except Exception as e:
        raise HTTPException(500, f"AI error: {e}")
